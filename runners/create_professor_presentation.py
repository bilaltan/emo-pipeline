#!/usr/bin/env python3
"""Generate an editable three-slide Phase 3.7/3.8 professor update deck."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = Path("results/phase37_scaling/phase37_scaling_professor_update.pptx")
NAVY = RGBColor(19, 42, 65)
TEAL = RGBColor(0, 128, 132)
GOLD = RGBColor(229, 171, 54)
PALE = RGBColor(232, 242, 242)
LIGHT = RGBColor(246, 248, 250)
GRAY = RGBColor(93, 104, 114)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, left, top, width, height, size=18, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_title(slide, title, subtitle):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    banner.fill.solid()
    banner.fill.fore_color.rgb = TEAL
    banner.line.fill.background()
    add_text(slide, title, 0.62, 0.34, 12.0, 0.5, size=29, bold=True)
    add_text(slide, subtitle, 0.64, 0.9, 12.0, 0.32, size=12, color=GRAY)
    add_text(slide, "EMO / Spark + Delta Lake", 9.75, 7.08, 2.85, 0.25, size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullet(slide, text, left, top, width, size=16, color=NAVY):
    box = add_text(slide, "", left, top, width, 0.45, size=size, color=color)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.level = 0
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.space_after = Pt(8)
    return box


def process_box(slide, label, detail, left, top, width, fill):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = TEAL
    add_text(slide, label, left + 0.1, top + 0.16, width - 0.2, 0.27, size=15, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, detail, left + 0.12, top + 0.52, width - 0.24, 0.27, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    return shape


def add_table(slide, rows, left, top, width, row_height=0.48):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, Inches(left), Inches(top), Inches(width), Inches(row_height * len(rows)))
    table = table_shape.table
    column_widths = [0.19, 0.2, 0.27, 0.15, 0.19]
    for index, fraction in enumerate(column_widths):
        table.columns[index].width = Inches(width * fraction)
    for row_index, values in enumerate(rows):
        for col_index, value in enumerate(values):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_index == 0 else (PALE if row_index % 2 else LIGHT)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER if col_index else PP_ALIGN.LEFT
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(11)
            paragraph.font.bold = row_index == 0
            paragraph.font.color.rgb = WHITE if row_index == 0 else NAVY
    return table_shape


def slide_one(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Full-Graph Learning Without One-Machine Graph Loading", "Phase 3.7/3.8: a general Spark + Delta execution path")
    add_text(slide, "Why this path?", 0.65, 1.42, 3.2, 0.32, size=18, bold=True, color=TEAL)
    add_bullet(slide, "Conventional GNNs repeat neighborhood aggregation every epoch.", 0.65, 1.84, 4.55, size=15)
    add_bullet(slide, "At Papers100M scale, repeated graph access creates high memory and I/O pressure.", 0.65, 2.48, 4.55, size=15)
    add_bullet(slide, "We compute graph features once, cache them in Delta, then train one shared classifier.", 0.65, 3.24, 4.55, size=15)
    process_box(slide, "Delta Graph Tables", "nodes + edges + splits", 5.5, 1.82, 1.63, PALE)
    process_box(slide, "Hop 1 Mean", "distributed Spark", 7.52, 1.82, 1.45, PALE)
    process_box(slide, "Hop 2 Mean", "cached to Delta", 9.35, 1.82, 1.45, PALE)
    process_box(slide, "Global Classifier", "one shared model", 11.18, 1.82, 1.52, RGBColor(240, 234, 211))
    for start, end in ((7.13, 7.52), (8.97, 9.35), (10.8, 11.18)):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(start), Inches(2.34), Inches(end), Inches(2.34))
        connector.line.color.rgb = TEAL
        connector.line.width = Pt(2)
    equation_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(3.25), Inches(7.2), Inches(1.15))
    equation_panel.fill.solid()
    equation_panel.fill.fore_color.rgb = LIGHT
    equation_panel.line.color.rgb = RGBColor(207, 216, 222)
    add_text(slide, "Node representation:  z = [x0 | x1 | x2]", 5.9, 3.53, 6.4, 0.28, size=21, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Original node features + one-hop mean + two-hop mean", 5.9, 3.92, 6.4, 0.23, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Key system property", 0.65, 4.72, 2.4, 0.3, size=18, bold=True, color=TEAL)
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(5.12), Inches(12.03), Inches(1.08))
    callout.fill.solid()
    callout.fill.fore_color.rgb = NAVY
    callout.line.fill.background()
    add_text(slide, "No driver-side graph collection. No unbounded Python adjacency lists. Spark processes bounded partitions; Delta provides reusable checkpoints.", 0.95, 5.43, 11.45, 0.38, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_two(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "One General Data Model, Three Graph Scales", "Same two-hop propagation and one global classifier across all experiments")
    rows = [
        ["Dataset", "Nodes", "Propagation edges", "Coverage", "Test accuracy"],
        ["WikiCS", "11,701", "431,206", "100%", "0.7746"],
        ["ogbn-products", "2,449,029", "123,718,024", "100%", "0.7068"],
        ["ogbn-papers100M", "111,059,956", "3,228,124,712", "100%", "0.6327"],
    ]
    add_table(slide, rows, 0.65, 1.55, 12.03)
    add_text(slide, "Experiment protocol", 0.65, 3.86, 3.0, 0.32, size=18, bold=True, color=TEAL)
    add_bullet(slide, "Two fixed neighborhood-mean hops, checkpointed in Delta Lake.", 0.65, 4.28, 5.8, size=15)
    add_bullet(slide, "Official OGB train/validation/test splits for Products and Papers100M.", 0.65, 4.88, 5.8, size=15)
    add_bullet(slide, "Executor count changes only system parallelism; model computation remains identical.", 0.65, 5.48, 5.8, size=15)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(4.05), Inches(5.53), Inches(1.85))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE
    panel.line.color.rgb = TEAL
    add_text(slide, "Research Interpretation", 7.5, 4.37, 4.85, 0.28, size=17, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Papers100M is the stress test, not a special-case implementation. Every dataset uses the same nodes-edges-splits abstraction.", 7.55, 4.78, 4.73, 0.73, size=15, color=NAVY, align=PP_ALIGN.CENTER)


def slide_three(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Scaling Helps Large Graphs, Then Saturates", "Phase 3.7 propagation time; accuracy and coverage are unchanged across executor counts")
    chart_data = CategoryChartData()
    chart_data.categories = ["8", "16", "32", "64"]
    chart_data.add_series("WikiCS", (41.1, 41.0, 45.5, 58.0))
    chart_data.add_series("Products", (110.4, 86.7, 81.7, 82.5))
    chart_data.add_series("Papers100M", (3199.1, 2459.0, 1912.0, 1890.7))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.65), Inches(1.45), Inches(7.0), Inches(4.65), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Executor instances"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Propagation time (seconds)"
    for series, color in zip(chart.series, (TEAL, GOLD, NAVY)):
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.5)
    add_text(slide, "Practical operating point", 8.08, 1.55, 4.2, 0.32, size=18, bold=True, color=TEAL)
    add_text(slide, "Papers100M", 8.08, 2.0, 2.2, 0.25, size=15, bold=True)
    add_text(slide, "8 -> 32 executors", 8.08, 2.32, 2.6, 0.25, size=13, color=GRAY)
    add_text(slide, "3,199.1 s -> 1,912.0 s", 8.08, 2.63, 4.2, 0.32, size=21, bold=True, color=NAVY)
    add_text(slide, "1.67x faster propagation", 8.08, 3.05, 3.8, 0.27, size=15, bold=True, color=TEAL)
    add_text(slide, "32 -> 64 executors saves only 21.3 seconds. Coordination and I/O overhead dominate beyond 32.", 8.08, 3.62, 4.25, 0.75, size=15, color=NAVY)
    takeaway = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.95), Inches(5.18), Inches(4.5), Inches(0.9))
    takeaway.fill.solid()
    takeaway.fill.fore_color.rgb = NAVY
    takeaway.line.fill.background()
    add_text(slide, "Recommendation: 32 executors for Papers100M", 8.15, 5.41, 4.1, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Next: compare this global path with bounded community GNNs and specialized distributed GNN systems using accuracy, runtime, memory, and cost.", 0.8, 6.38, 11.8, 0.42, size=14, color=GRAY, align=PP_ALIGN.CENTER)


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide_one(presentation)
    slide_two(presentation)
    slide_three(presentation)
    presentation.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
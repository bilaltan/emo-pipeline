#!/usr/bin/env python3
"""Generate a simple editable three-slide professor update deck."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = Path("results/phase37_scaling/phase37_scaling_professor_update.pptx")
BLUE = RGBColor(31, 78, 121)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(221, 235, 247)


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.55))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = BLUE


def add_bullets(slide, bullets, top=1.25, size=20):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(5.8))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BLACK
        paragraph.space_after = Pt(16)


def add_table(slide, rows, top, widths):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.6), Inches(top), Inches(12.1), Inches(0.55 * len(rows)))
    table = table_shape.table
    for col, width in enumerate(widths):
        table.columns[col].width = Inches(width)
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if row_index == 0 else (LIGHT_BLUE if row_index % 2 else WHITE)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(12)
            paragraph.font.bold = row_index == 0
            paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK
            paragraph.alignment = PP_ALIGN.CENTER


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "New Global Graph Learning Path")
    add_bullets(slide, [
        "The old path trains separate bounded GNNs for graph communities.",
        "The new path stores nodes, edges, and official splits in Delta Lake.",
        "Spark computes two neighbor-mean propagation hops across the distributed graph.",
        "The propagated features are cached in Delta Lake.",
        "One global logistic-regression classifier is trained on all training nodes.",
        "No full graph or large adjacency list is collected on the driver.",
    ], size=19)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Experiment Coverage")
    add_table(slide, [
        ["Dataset", "Nodes", "Propagation edges", "Node coverage", "Test accuracy"],
        ["WikiCS", "11,701", "431,206", "100%", "0.7746"],
        ["ogbn-products", "2,449,029", "123,718,024", "100%", "0.7068"],
        ["ogbn-papers100M", "111,059,956", "3,228,124,712", "100%", "0.6327"],
    ], 1.4, [2.0, 2.0, 3.1, 2.2, 2.8])
    add_bullets(slide, [
        "The same nodes-edges-splits pipeline was used across all three graph sizes.",
        "Executor count changes runtime only; accuracy is unchanged because the computation is the same.",
        "Papers100M is the large-scale stress test of the general architecture.",
    ], top=4.25, size=19)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Executor Scaling Results")
    add_table(slide, [
        ["Dataset", "8 executors", "Best propagation time", "Recommended executors"],
        ["WikiCS", "41.1 s", "41.0 s at 16", "8"],
        ["ogbn-products", "110.4 s", "81.7 s at 32", "32"],
        ["ogbn-papers100M", "3,199.1 s", "1,912.0 s at 32", "32"],
    ], 1.4, [2.4, 2.4, 4.0, 3.3])
    add_bullets(slide, [
        "Papers100M propagation improved 1.67x from 8 to 32 executors.",
        "Increasing from 32 to 64 executors gave only a small additional improvement.",
        "Next: compare this path with the community GNN path and specialized distributed GNN baselines.",
    ], top=4.25, size=19)

    presentation.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()